#!/home/wnstn4647/AI-CHAT/.venv/bin/python3
"""
채널A AI 문서검색시스템 PowerPoint 생성 스크립트
Markdown 파일을 읽어서 실제 .pptx 파일 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    """PowerPoint 프레젠테이션 생성"""

    # 새 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 슬라이드 1: 표지
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "채널A AI 문서 검색 시스템"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "RAG 기반 지능형 문서 검색 플랫폼"
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 발표 정보
    info_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
    info_frame = info_box.text_frame
    info_frame.text = "발표일: 2024.11.24\n발표자: [이름]"
    info_frame.paragraphs[0].font.size = Pt(18)
    info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 슬라이드 2: 목차
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "목차"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. 시스템 개요 - 무엇을, 왜 만들었나?"

    for item in [
        "2. 주요 기능 - 어떤 일을 하나?",
        "3. 기술 아키텍처 - 어떻게 작동하나?",
        "4. 검색 메커니즘 - 어떻게 찾나?",
        "5. 실제 사용 예시 - 어떻게 쓰나?",
        "6. 운영 현황 - 지금 어떤 상태인가?",
        "7. 기술적 특징 - 무엇이 특별한가?",
        "8. 향후 계획 - 앞으로 어떻게 할 것인가?"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

    # 슬라이드 3: 시스템 개요
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "시스템 개요"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "문제 상황"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = '"작년에 작성한 중계차 장비 구매 문서, 어디 있더라...?"'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '"폴더 10개를 뒤져서 찾은 파일이 잘못된 버전..."'
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "해결 방법: AI 문서 검색 시스템"
    p.font.bold = True
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = "📚 475개 문서 (2014~2025, 12년치)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "⚡ 초고속 검색 (BM25 + FTS5)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "🤖 AI가 내용까지 요약"
    p.level = 1

    # 슬라이드 4: 주요 기능 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "주요 기능 (1/2)"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1️⃣ SEARCH 모드 - 문서 찾기"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = '사용자: "중계차 카메라 문서 찾아줘"'
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ 빠르게 관련 문서 3~5개 제시"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "(기안자, 날짜, 금액 정보 포함)"
    p.level = 2

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "2️⃣ DOCUMENT 모드 - 문서 요약"
    p.font.bold = True
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = '사용자: "이 문서 요약해줘"'
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ AI가 2~3초 만에 핵심 내용 요약"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "(제목, 날짜, 금액, 주요 내용)"
    p.level = 2

    # 슬라이드 5: 주요 기능 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "주요 기능 (2/2)"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "3️⃣ QA 모드 - 질문/답변"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = '사용자: "렌즈 오버홀 비용은 얼마였어?"'
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ AI가 관련 문서 찾아서 답변 생성"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "(3~5초, 자연어 답변)"
    p.level = 2

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "4️⃣ COST 모드 - 비용 조회"
    p.font.bold = True
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = '사용자: "2024년 남준수 기안 문서 총액은?"'
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ SQL 직접 집계 (초고속)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "(초고속 비용 합계)"
    p.level = 2

    # 슬라이드 6: 기술 아키텍처 (초보자용)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "기술 아키텍처 (초보자용)"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "시스템 = 디지털 도서관"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "📚 PDF 파일 = 책"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "📇 metadata.db = 도서 카드 목록"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "📑 BM25 인덱스 = 책 뒤의 색인(index)"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "🤖 AI(LLM) = 사서 (내용 요약해주는)"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "🔍 검색 엔진 = 도서 검색 시스템"
    p.level = 0

    # 슬라이드 7: 기술 아키텍처 (전문가용)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "기술 아키텍처 (전문가용)"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Dual-Server Architecture"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Streamlit (UI 서버) - Port 8501"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "FastAPI (API 서버) - Port 7860"
    p.level = 0

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "RAG Pipeline (통합 파사드)"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "→ Query Router (쿼리 분류)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ Hybrid Search (BM25 + FTS5)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ LLM (Qwen 2.5, llama.cpp)"
    p.level = 1

    # 슬라이드 8: 검색 메커니즘
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "검색 메커니즘: 3단계 하이브리드"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Stage 1: ExactMatch (정확한 파일명)"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "→ 파일명에서 직접 매칭"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Stage 2: FTS5 (전문검색)"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "→ SQLite Full-Text Search"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Stage 3: BM25 (키워드 랭킹)"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "→ TF-IDF 기반 순위 알고리즘"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ 가장 관련성 높은 문서 반환"
    p.level = 1

    # 슬라이드 9: 실제 사용 예시
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "실제 사용 예시"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "시나리오 1: 급하게 문서 찾기"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = '질문: "중계차 카메라 렌즈 관련 문서"'
    p.level = 1

    p = tf.add_paragraph()
    p.text = "결과: 3개 문서 즉시 발견"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "시나리오 2: 비용 확인"
    p.font.bold = True
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = '질문: "2024년 장비 구매 총액?"'
    p.level = 1

    p = tf.add_paragraph()
    p.text = "결과: 즉시 금액 합계 계산"
    p.level = 1

    # 슬라이드 10: 운영 현황
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "운영 현황 (2024.11.24 기준)"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "📊 시스템 통계"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "문서 수: 475개 PDF"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "기간: 2014 ~ 2025 (12년)"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "텍스트 추출: 87% (나머지는 OCR 처리 가능)"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "자동 Fallback OCR 지원"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Dual-Server 아키텍처 (현재 실행 중)"
    p.level = 0

    # 슬라이드 11: 기술적 특징
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "기술적 특징"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. 한국어 특화"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "형태소 분석, 동의어 처리"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "2. OCR Fallback 처리"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "텍스트 부족 시 자동 OCR 실행"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "3. 3단계 캐싱"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "메모리 → 디스크 → DB"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "4. 로컬 실행"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "외부 유출 없음 (데이터 보안)"
    p.level = 1

    # 슬라이드 12: 향후 계획
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "향후 계획"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "단기 (1~3개월)"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "문서 자동 업로드 시스템"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "사용자 피드백 수집"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "중기 (3~6개월)"
    p.font.bold = True
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "문서 10,000개로 확장"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "다른 부서 도입"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "장기 (6개월 이상)"
    p.font.bold = True
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "회사 전체 통합 검색 시스템"
    p.level = 1

    # 슬라이드 13: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # 중앙에 큰 텍스트
    qa_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    qa_frame = qa_box.text_frame
    qa_frame.text = "Q & A\n\n질문 있으시면 말씀해주세요"
    qa_frame.paragraphs[0].font.size = Pt(48)
    qa_frame.paragraphs[0].font.bold = True
    qa_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    qa_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return prs


def main():
    """메인 함수"""

    # 출력 경로
    output_dir = "/home/wnstn4647/AI-CHAT/docs/발표자료"
    output_path = os.path.join(output_dir, "채널A_AI문서검색시스템.pptx")

    print("🎬 PowerPoint 생성 시작...")

    # 프레젠테이션 생성
    prs = create_presentation()

    # 저장
    prs.save(output_path)

    print(f"✅ PowerPoint 생성 완료!")
    print(f"📁 파일 위치: {output_path}")
    print(f"📊 슬라이드 수: {len(prs.slides)}")

    # 파일 정보
    file_size = os.path.getsize(output_path)
    print(f"📦 파일 크기: {file_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
